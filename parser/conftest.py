"""Shared fixtures for the parser test suite.

Living at the package root puts `parser/` on sys.path so tests can
`import main` (also enforced by `pythonpath = .` in pytest.ini).
"""
import io

import pytest
from docx import Document
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfgen import canvas

from main import app


@pytest.fixture
def client() -> TestClient:
    return TestClient(app)


@pytest.fixture
def pdf_bytes() -> bytes:
    """A valid one-page PDF containing the text 'hello pdf'."""
    buf = io.BytesIO()
    c = canvas.Canvas(buf)
    c.drawString(72, 720, "hello pdf")
    c.showPage()
    c.save()
    return buf.getvalue()


@pytest.fixture
def docx_bytes() -> bytes:
    """A valid DOCX containing the text 'hello docx'."""
    doc = Document()
    doc.add_paragraph("hello docx")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


@pytest.fixture
def pptx_bytes() -> bytes:
    """A valid one-slide PPTX containing the text 'hello pptx'."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "hello pptx"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
