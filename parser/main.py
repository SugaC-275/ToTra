# parser/main.py
"""ToTra document parser — extracts plain text from PDF, DOCX, and PPTX uploads."""
import io
import logging
import os
import time
from collections.abc import Callable

import pdfplumber
from docx import Document
from fastapi import FastAPI, File, HTTPException, UploadFile
from pptx import Presentation

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(name)s %(message)s",
)
logger = logging.getLogger("totra.parser")

# Maximum accepted upload size in bytes (default 10 MiB).
MAX_UPLOAD_BYTES = int(os.getenv("MAX_UPLOAD_BYTES", str(10 * 1024 * 1024)))

app = FastAPI(title="ToTra Parser")


def _elapsed_ms(start: float) -> int:
    return int((time.monotonic() - start) * 1000)


def _parse_pdf(data: bytes) -> tuple[str, int]:
    parts: list[str] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        page_count = len(pdf.pages)
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                parts.append(text)
    return "\n".join(parts), page_count


def _parse_docx(data: bytes) -> tuple[str, int]:
    doc = Document(io.BytesIO(data))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    # DOCX has no reliable page count without rendering; report 1. See README.
    return text, 1


def _parse_pptx(data: bytes) -> tuple[str, int]:
    prs = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for slide in prs.slides:
        texts = [
            shape.text
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        slides.append("\n".join(texts))
    return "\n\n".join(slides), len(prs.slides)


_PARSERS: dict[str, Callable[[bytes], tuple[str, int]]] = {
    "pdf": _parse_pdf,
    "docx": _parse_docx,
    "pptx": _parse_pptx,
}


def _extract(filename: str, data: bytes, size: int) -> tuple[str, str, int]:
    """Validate the upload and dispatch to the right parser.

    Returns (format, text, page_count). Raises HTTPException for any bad input.
    """
    if size == 0:
        raise HTTPException(status_code=400, detail="file is empty")
    if size > MAX_UPLOAD_BYTES:
        raise HTTPException(
            status_code=413,
            detail=f"file exceeds the {MAX_UPLOAD_BYTES}-byte limit",
        )

    fmt = filename.rpartition(".")[2]
    parser = _PARSERS.get(fmt)
    if parser is None:
        raise HTTPException(status_code=400, detail="unsupported format")

    try:
        text, page_count = parser(data)
    except Exception as exc:  # any parser failure means the file is unreadable
        logger.warning("parse failed file=%r: %s", filename, exc)
        raise HTTPException(
            status_code=422, detail=f"could not parse {fmt} file"
        ) from exc
    return fmt, text, page_count


@app.get("/health")
async def health():
    return {"status": "ok"}


@app.post("/parse")
async def parse_file(file: UploadFile = File(...)):
    start = time.monotonic()
    filename = (file.filename or "").lower()
    data = await file.read()
    size = len(data)

    try:
        fmt, text, page_count = _extract(filename, data, size)
    except HTTPException as exc:
        logger.info(
            "parse rejected file=%r size=%d status=%d duration_ms=%d",
            filename, size, exc.status_code, _elapsed_ms(start),
        )
        raise

    logger.info(
        "parse ok file=%r size=%d format=%s page_count=%d duration_ms=%d",
        filename, size, fmt, page_count, _elapsed_ms(start),
    )
    return {"text": text, "page_count": page_count}
